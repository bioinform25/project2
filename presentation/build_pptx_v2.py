#!/usr/bin/env python3
"""Rebuild the dry-lab summary deck to match the visual style of
'260713 lab meeting.pptx' (VS-Code-dark code boxes, Arial titles,
white-bordered result boxes/tables, Courier/Calibri results)."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

BASE = r"C:\Users\SAMSUNG\Desktop\project2"
FIG = os.path.join(BASE, "figures")
OUT = os.path.join(BASE, "presentation", "dry_lab_summary.pptx")

# ---- palette (matched to the reference template) ----
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
CODE_BORDER = RGBColor(0x45, 0x45, 0x45)
GREEN_COMMENT = RGBColor(0x6A, 0x99, 0x55)
YELLOW_KEY = RGBColor(0xDC, 0xDC, 0xAA)
GRAY_CODE = RGBColor(0xD4, 0xD4, 0xD4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def EI(inches):
    """Inches -> integer EMU, safe against float division bugs."""
    return Emu(int(Inches(inches)))


def new_slide():
    return prs.slides.add_slide(BLANK)


def set_run(run, text, size, bold=False, italic=False, color=BLACK, font="Arial"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def add_content_title(slide, text):
    """Top title bar matching the reference: (0.4,0.15,12.5,0.55), Arial 20 bold black."""
    tb = slide.shapes.add_textbox(EI(0.4), EI(0.15), EI(12.5), EI(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, text, 20, bold=True)
    return tb


def add_section_header(title, subtitle):
    """Divider slide matching the reference 'PROJECT N' slides: big bold title +
    smaller subtitle, centered around y=3.0."""
    s = new_slide()
    tb = s.shapes.add_textbox(EI(0.8), EI(3.0), EI(11.7), EI(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    set_run(r0, title, 32, bold=True)
    p1 = tf.add_paragraph()
    p1.space_before = Pt(10)
    r1 = p1.add_run()
    set_run(r1, subtitle, 16)
    return s


def add_overview_slide(title, lines):
    """Title + paragraph body matching reference slide 2 (Overall Study Flow):
    title bar, then a body textbox of '- ' prefixed lines."""
    s = new_slide()
    add_content_title(s, title)
    tb = s.shapes.add_textbox(EI(0.6), EI(1.2), EI(12.1), EI(5.9))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line == "":
            p.space_before = Pt(8)
            r = p.add_run()
            set_run(r, "", 14)
            continue
        r = p.add_run()
        set_run(r, line, 14)
        p.line_spacing = 1.2
    return s


def add_code_box(slide, left_in, top_in, width_in, height_in, lines):
    """lines: list of (text, kind) with kind in {'comment','key','code'}.
    Matches reference: rounded rect, fill 1E1E1E, border 454545, Consolas 7pt,
    108% line spacing, vertically centered; comment lines green, key lines
    yellow, regular code gray; first line centered if it's a comment header."""
    color_map = {"comment": GREEN_COMMENT, "key": YELLOW_KEY, "code": GRAY_CODE}
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  EI(left_in), EI(top_in), EI(width_in), EI(height_in))
    box.adjustments[0] = 0.04
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = CODE_BORDER
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (text, kind) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.08
        p.space_after = Pt(0)
        if i == 0 and kind == "comment":
            p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, text, 7, color=color_map[kind], font="Consolas")
    return box


def add_explanation(slide, left_in, top_in, width_in, height_in, text):
    """Plain textbox, no border: Arial 11.5, matching reference explanation box."""
    tb = slide.shapes.add_textbox(EI(left_in), EI(top_in), EI(width_in), EI(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        r = p.add_run()
        set_run(r, line, 11.5)
    return tb


def add_result_table(slide, left_in, top_in, width_in, height_in, headers, rows, font_size=10):
    """White bg, bold Calibri header, black border - matching reference tables."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gshape = slide.shapes.add_table(n_rows, n_cols, EI(left_in), EI(top_in),
                                     EI(width_in), EI(height_in))
    table = gshape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
        r = cell.text_frame.paragraphs[0].add_run()
        set_run(r, str(h), font_size, bold=True, font="Calibri")
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            r = cell.text_frame.paragraphs[0].add_run()
            set_run(r, str(val), font_size, font="Calibri")
    return gshape


def add_result_image(slide, left_in, top_in, width_in, height_in, img_path):
    max_w = EI(width_in)
    max_h = EI(height_in)
    pic = slide.shapes.add_picture(img_path, EI(left_in), EI(top_in), height=max_h)
    if pic.width > max_w:
        ratio = max_w / pic.width
        pic.width = int(max_w)
        pic.height = int(pic.height * ratio)
    pic.left = int(EI(left_in) + (max_w - pic.width) / 2)
    pic.top = int(EI(top_in) + (max_h - pic.height) / 2)
    return pic


def content_slide(title, code_lines, explanation, result_image=None, result_table=None):
    """Full layout matching the reference content slides."""
    s = new_slide()
    add_content_title(s, title)
    code_h = 4.0 if (result_table or result_image) else 4.0
    add_code_box(s, 0.4, 1.05, 7.35, 4.15, code_lines)
    add_explanation(s, 7.95, 1.05, 4.98, 1.3, explanation)
    if result_image:
        add_result_image(s, 7.95, 2.55, 4.98, 3.65, result_image)
    elif result_table:
        headers, rows = result_table
        row_h = 3.65 / (len(rows) + 1)
        add_result_table(s, 7.95, 2.55, 4.98, 3.65, headers, rows)
    return s


# =====================================================================
# SLIDE 1: TITLE
# =====================================================================
s = new_slide()
tb = s.shapes.add_textbox(EI(1.0), EI(2.2), EI(11.3), EI(3.5))
tf = tb.text_frame
tf.word_wrap = True
p0 = tf.paragraphs[0]
r0 = p0.add_run()
set_run(r0, "PHMG → HSF2 → NETosis Hypothesis: Dry-Lab Verification", 30, bold=True)
p1 = tf.add_paragraph()
p1.space_before = Pt(16)
r1 = p1.add_run()
set_run(r1, "공개 human 섬유화 코호트(GSE135251, GSE14323, GSE47460, GSE53845) + "
            "HCC 면역항암제 코호트(GSE215011, GSE279750, GSE235863) 분석", 18)
p2 = tf.add_paragraph()
p2.space_before = Pt(10)
r2 = p2.add_run()
set_run(r2, "github.com/bioinform25/project2", 13)

# =====================================================================
# SLIDE 2: OVERVIEW
# =====================================================================
add_overview_slide("Overall Study Flow", [
    "- 연구실 웻랩 실험: PHMG 6 mg/kg in vivo 마우스 호중구 RNA-seq → 상승(upregulated) 유전자군의 "
    "1순위 상위 TF로 HSF2 확인; ex vivo PHMG 3 ug/mL, 6시간 처리 → Western blot에서 HSF2 band "
    "두꺼워짐 관찰.",
    "- 작업 가설: PHMG → HSF2 활성화 → N2(면역억제성) 호중구 극성화 → NETosis 유전자 발현 유도 → "
    "섬유화 / 간암(HCC) 면역회피.",
    "- 이 자료의 목적: 추가 웻랩 실험 전에, 이 가설의 population-level 예측(HSF2↑ → NETosis 핵심 "
    "유전자↑)을 4개의 독립적인 공개 human 섬유화 코호트로 먼저 검증하고, 3개의 독립적인 HCC 면역항암제 "
    "반응성 코호트로 확장 분석.",
    "- 총 12개 R 스크립트(repo의 scripts/00-12) — 결과를 보기 전에 미리 지정한 분석계획: 1차 endpoint, "
    "재현성 규칙, 특이성 대조, in silico TF 검증, 전체 meta-analysis까지 포함.",
    "- 이 PPT는 실제 분석이 진행된 순서 그대로, 각 스크립트의 코드·실제 결과·한 줄 해석을 정리했습니다.",
])

# =====================================================================
# SLIDE 3: SECTION HEADER - PART 1
# =====================================================================
add_section_header(
    "PART 1 — Human Fibrosis Cohorts",
    "섬유화된 간·폐 조직에서 HSF2가 NETosis 핵심 유전자와 양의 상관관계를 보이는가?"
)

# =====================================================================
# SLIDE 4: PRE-SPECIFIED PLAN
# =====================================================================
code = [
    ("# scripts/00_functions.R - plan written BEFORE any result was seen", "comment"),
    ("NET_CORE     <- c(\"PADI4\", \"ELANE\", \"MPO\")", "key"),
    ("NET_EXTENDED <- c(NET_CORE, \"CTSG\", \"LTF\", \"CAMP\", \"DEFA4\",", "code"),
    ("                   \"AZU1\", \"ITGAM\", \"NCF2\", \"S100A12\", \"LCN2\")", "code"),
    ("TARGET_TF    <- \"HSF2\"", "key"),
    ("", "code"),
    ("# Primary endpoint: Spearman rho(HSF2, NET_core_score)", "comment"),
    ("#   within diseased/fibrotic samples, alpha = 0.05", "comment"),
    ("# Replication rule: primary + replication cohort must agree", "comment"),
    ("#   in sign AND both be nominally significant (p<0.05)", "comment"),
]
content_slide(
    "00_functions.R — Pre-Specified Hypothesis & Analysis Plan",
    code,
    "HSF2가 NETosis 유전자 발현을 유도한다면, 섬유화 조직에서 HSF2는 PADI4/ELANE/MPO와 양의 "
    "상관관계를 보여야 하고, 독립적인 코호트들에서 재현되어야 함. 이 계획은 실제 상관관계를 "
    "계산해보기 전에 미리 고정함.",
    result_table=(
        ["Cohort", "Organ", "Role", "n (fibrotic)"],
        [
            ["GSE135251", "Liver", "Primary", "206"],
            ["GSE14323", "Liver", "Replication", "96"],
            ["GSE47460", "Lung", "Primary", "122"],
            ["GSE53845", "Lung", "Replication", "40"],
        ],
    ),
)

# =====================================================================
# SLIDE 5: GSE135251 (liver, primary)
# =====================================================================
code = [
    ("# scripts/01_liver_GSE135251.R (RNA-seq, n=206 NAFLD)", "comment"),
    ("net_core_score <- panel_score(expr_panel, NET_CORE, min_genes = 2)", "code"),
    ("df_nafld <- df[df$disease == \"NAFLD\", ]", "code"),
    ("", "code"),
    ("results$primary <- spearman_report(", "key"),
    ("    df_nafld$HSF2, df_nafld$NET_core,", "key"),
    ("    \"HSF2 vs NET_core (NAFLD subset, primary)\")", "key"),
]
content_slide(
    "01_liver_GSE135251.R — Liver, Primary Cohort (n=206)",
    code,
    "결과는 원 가설과 반대 방향: 섬유화된 간 조직에서 HSF2는 NETosis_core score와 음의 상관관계를 "
    "보임.",
    result_image=os.path.join(FIG, "liver_GSE135251_HSF2_vs_NETcore.png"),
)

# =====================================================================
# SLIDE 6: GSE14323 (liver, replication)
# =====================================================================
code = [
    ("# scripts/02_liver_GSE14323.R (Affymetrix, independent platform)", "comment"),
    ("df_disease <- df[df$tissue != \"Normal\", ]  # n=96", "code"),
    ("", "code"),
    ("results$primary <- spearman_report(", "key"),
    ("    df_disease$HSF2, df_disease$NET_core,", "key"),
    ("    \"HSF2 vs NET_core (diseased subset, replication)\")", "key"),
    ("", "code"),
    ("# Replication rule (same sign + both p<0.05): MET", "comment"),
]
content_slide(
    "02_liver_GSE14323.R — Liver, Replication Cohort (n=96)",
    code,
    "독립적인 플랫폼, 독립적인 환자군 — 같은 음의 방향, 훨씬 강함(rho=-0.550, p=6.2e-09). 간 코호트 "
    "결과가 공식적으로 재현(REPLICATE)됨.",
    result_image=os.path.join(FIG, "liver_GSE14323_HSF2_vs_NETcore.png"),
)

# =====================================================================
# SLIDE 7: GSE47460 (lung, primary)
# =====================================================================
code = [
    ("# scripts/03_lung_GSE47460.R (Agilent, n=122 UIP/IPF)", "comment"),
    ("df_fibrotic <- df[df$group == \"UIP_IPF\", ]", "code"),
    ("", "code"),
    ("results$primary <- spearman_report(", "key"),
    ("    df_fibrotic$HSF2, df_fibrotic$NET_core,", "key"),
    ("    \"HSF2 vs NET_core (UIP/IPF subset, primary)\")", "key"),
]
content_slide(
    "03_lung_GSE47460.R — Lung, Primary Cohort (n=122)",
    code,
    "다른 장기(간이 아닌 폐), 다른 플랫폼 — 같은 음의 방향, 매우 유의함(rho=-0.392, p=8.1e-06). "
    "간에만 국한된 패턴이 아님.",
    result_image=os.path.join(FIG, "lung_GSE47460_HSF2_vs_NETcore.png"),
)

# =====================================================================
# SLIDE 8: GSE53845 (lung, replication)
# =====================================================================
code = [
    ("# scripts/04_lung_GSE53845.R (Agilent 2-color, n=40 IPF)", "comment"),
    ("df_ipf <- df[df$diagnosis == \"IPF\", ]", "code"),
    ("", "code"),
    ("results$primary <- spearman_report(", "key"),
    ("    df_ipf$HSF2, df_ipf$NET_core,", "key"),
    ("    \"HSF2 vs NET_core (IPF subset, replication)\")", "key"),
    ("", "code"),
    ("# rho=-0.050, p=0.760 (n.s.) - same sign, underpowered (n=40)", "comment"),
]
content_slide(
    "04_lung_GSE53845.R — Lung, Replication Cohort (n=40)",
    code,
    "방향은 여전히 음의 상관이지만 n=40은 검정력 부족 — 개별적으로는 비유의. \"반증\"이 아니라 "
    "\"정보 없음\"으로 취급하고, meta-analysis에 그대로 포함시킴.",
    result_image=os.path.join(FIG, "lung_GSE53845_HSF2_vs_NETcore.png"),
)

# =====================================================================
# SLIDE 9: META-ANALYSIS (KEY FINDING)
# =====================================================================
code = [
    ("# scripts/05_meta_analysis.R", "comment"),
    ("# from-scratch Fisher-z / DerSimonian-Laird random-effects", "comment"),
    ("meta_all <- meta_cor_DL(primary_df$rho, primary_df$n, primary_df$cohort)", "key"),
    ("", "code"),
    ("cat(sprintf(\"pooled rho=%.3f, 95%% CI [%.3f,%.3f], p=%.3g", "code"),
    ("Heterogeneity: I^2=%.1f%%\", meta_all$pooled_rho, meta_all$ci_lo,", "code"),
    ("meta_all$ci_hi, meta_all$p_value, meta_all$I2))", "code"),
]
content_slide(
    "05_meta_analysis.R — KEY FINDING: 4-Cohort Meta-Analysis",
    code,
    "Pooled rho=-0.308(p=0.010)는 유의하며 가설과 반대 방향 — GSE53845 단독은 검정력 부족이었지만 "
    "무관함. 다만 I^2=83.5%(높은 이질성): 유의성은 \"방향이 진짜\"라는 뜻이고, 이질성은 \"효과 크기가 "
    "코호트마다 일정하지 않다\"는 별개의 질문. 이 크기 문제를 해결해줄 게 웻랩 time-course 실험.",
    result_image=os.path.join(FIG, "FOREST_HSF2_vs_NETcore_all_cohorts.png"),
)

# =====================================================================
# SLIDE 10: SPECIFICITY CONTROL
# =====================================================================
code = [
    ("# scripts/06_specificity_check_PTPRC.R", "comment"),
    ("# Is this just \"more immune infiltrate dilutes HSF2 signal\"?", "comment"),
    ("# Test against PTPRC (CD45, pan-leukocyte marker):", "comment"),
    ("r1_ptprc   <- spearman_report(df1_dis$HSF2, df1_dis$PTPRC, ...)", "key"),
    ("r1_netcore <- spearman_report(df1_dis$HSF2, df1_dis$NET_core, ...)", "key"),
]
content_slide(
    "06_specificity_check_PTPRC.R — Ruling Out a Trivial Artifact",
    code,
    "HSF2는 전체 면역세포 침윤(PTPRC)과는 함께 증가(양의 상관)하지만 NETosis 핵심 유전자와는 "
    "특이적으로 감소(음의 상관) — 서로 반대 방향. \"면역세포가 많아져서 HSF2 상대신호가 희석됐다\"는 "
    "단순한 설명을 배제함.",
    result_table=(
        ["Cohort", "HSF2 vs PTPRC", "HSF2 vs NET_core"],
        [
            ["GSE135251", "+0.178 (p=0.010)", "-0.145 (p=0.038)"],
            ["GSE14323", "+0.448 (p=4.7e-06)", "-0.550 (p=6.2e-09)"],
        ],
    ),
)

# =====================================================================
# SLIDE 11: CROSS-COHORT GENE CONSISTENCY
# =====================================================================
code = [
    ("# scripts/08_cross_cohort_gene_consistency.R", "comment"),
    ("# Which individual NETosis genes move with HSF2 the SAME", "comment"),
    ("# direction in all 4 cohorts, vs. which flip liver<->lung?", "comment"),
    ("consistency$n_cohorts_sig_FDR05 <- rowSums(consistency[, fdr_cols] < 0.05)", "code"),
    ("consistency$direction_consistent <-", "key"),
    ("    (consistency$n_negative == consistency$n_cohorts_tested) |", "key"),
    ("    (consistency$n_positive == consistency$n_cohorts_tested)", "key"),
]
content_slide(
    "08_cross_cohort_gene_consistency.R — Which Genes Drive This",
    code,
    "CAMP/MPO/DEFA4(+ 작은 예외 하나씩 있는 PADI4/ELANE)는 4개 코호트 전부에서 HSF2와 일관되게 "
    "음의 방향. LTF/NCF2/ITGAM/LCN2는 간과 폐 사이에서 부호가 뒤집힘 — 호중구 특이적 신호가 아닐 "
    "가능성(예: LCN2는 스트레스 받은 간세포에서도 만들어짐).",
    result_table=(
        ["Gene", "GSE135251", "GSE14323", "GSE47460", "GSE53845", "Direction"],
        [
            ["CAMP", "-0.056", "-0.459", "-0.457", "-0.529", "Consistent (-)"],
            ["MPO", "-0.025", "-0.277", "-0.320", "-0.412", "Consistent (-)"],
            ["DEFA4", "-0.090", "-0.337", "-0.328", "-0.097", "Consistent (-)"],
            ["PADI4", "-0.120", "-0.496", "-0.205", "+0.181", "3 neg/1 pos (ns)"],
            ["ELANE", "-0.192", "-0.422", "-0.348", "+0.075", "3 neg/1 pos (ns)"],
        ],
    ),
)

# =====================================================================
# SLIDE 12: CHEA3 IN SILICO CHECK
# =====================================================================
code = [
    ("# scripts/09_chea3_analysis.R", "comment"),
    ("# Does any existing ChIP-seq/co-expression DB already show", "comment"),
    ("# HSF2 upstream of these genes?", "comment"),
    ("res <- httr::POST(", "key"),
    ("  url = \"https://maayanlab.cloud/chea3/api/enrich/\",", "key"),
    ("  body = list(query_name = query_name, gene_set = as.list(genes)))", "key"),
]
content_slide(
    "09_chea3_analysis.R — No Existing Evidence for This Regulation",
    code,
    "8개의 증거 라이브러리(ChIP-seq + co-expression + literature) 전체에서 HSF2는 이 유전자 "
    "세트의 상위 조절인자로 한 번도 나오지 않음. 이 도구들이 못 잡는 억제성 관계이거나, 정말 "
    "미개척 영역(blue ocean)일 가능성.",
    result_table=(
        ["Library", "HSF2 rank", "Total TFs", "Overlap"],
        [
            ["Integrated mean rank", "1444", "1632", "-"],
            ["ARCHS4 co-expression", "1209", "1628", "0"],
            ["GTEx co-expression", "1417", "1607", "0"],
            ["ENCODE ChIP-seq", "not returned", "118", "-"],
            ["ReMap ChIP-seq", "not returned", "297", "-"],
        ],
    ),
)

# =====================================================================
# SLIDE 13: SECTION HEADER - PART 2
# =====================================================================
add_section_header(
    "PART 2 — HCC Immunotherapy Cohorts",
    "HCC = 간암, ICI = 면역관문억제제. 실제 환자의 종양 HSF2가 면역항암제 반응 여부와 관련 있는가?"
)

# =====================================================================
# SLIDE 14: GATE-CHECK + GSE215011
# =====================================================================
code = [
    ("# GSE140901 (originally planned dataset) - GATE-CHECK first:", "comment"),
    ("# 785-gene NanoString panel - does it measure our genes?", "comment"),
    ("zcat GSE140901_processed_data.txt.gz | grep -w \"HSF2\"", "key"),
    ("# >>> no output - HSF2 NOT in panel; PADI4/MPO also missing", "comment"),
    ("# ==> EXCLUDED, not force-fit with a substitute readout", "comment"),
    ("", "code"),
    ("# scripts/07_hcc_ICI_GSE215011.R (nivolumab monotherapy, n=10)", "comment"),
    ("wilcox_report(df$HSF2, df$group, \"HSF2: Responder vs Non-responder\")", "key"),
]
content_slide(
    "Dataset Gate-Check + First Cohort: GSE215011 (n=10)",
    code,
    "원래 계획했던 데이터셋은 HSF2를 아예 측정하지 않음 — 직접 확인 후 대체하지 않고 배제. 대신 "
    "GSE215011(실제 RNA-seq)을 찾아 사용. n=10 단독으로는 유의성 도달에 표본이 부족함"
    "(rank-biserial r=0.44, p=0.30).",
    result_image=os.path.join(FIG, "hcc_ICI_GSE215011_HSF2_vs_NETcore.png"),
)

# =====================================================================
# SLIDE 15: GSE279750 + GSE235863
# =====================================================================
code = [
    ("# scripts/10_hcc_ICI_GSE279750.R (anti-PD-L1 combo, post-tx, n=10)", "comment"),
    ("wilcox_report(df$HSF2, df$group, \"HSF2: Responder vs Non-responder\")", "key"),
    ("# >>> rank-biserial r=-0.92, p=0.025 (responders have LOWER HSF2)", "comment"),
    ("", "code"),
    ("# scripts/11_hcc_ICI_GSE235863.R (anti-PD-1+lenvatinib, n=15)", "comment"),
    ("wilcox_report(df$HSF2, df$group, \"HSF2: Responder vs Non-responder\")", "key"),
    ("# >>> rank-biserial r=-0.73, p=0.043 (same direction, also sig.)", "comment"),
]
content_slide(
    "Two More Independent HCC-ICI Cohorts (n=10, n=15)",
    code,
    "두 코호트 모두에서 종양 HSF2가 낮은 환자일수록 병용 면역항암제에 더 잘 반응(둘 다 p<0.05). "
    "참고: 둘 다 병용요법·치료 후 코호트 — GSE215011의 단독요법과는 설계가 다름.",
    result_image=os.path.join(FIG, "hcc_ICI_GSE279750_HSF2_by_response.png"),
)

# =====================================================================
# SLIDE 16: ICI POOLED META-ANALYSIS
# =====================================================================
code = [
    ("# scripts/12_hcc_ICI_pooled_summary.R", "comment"),
    ("# rank-biserial r = scale-free Wilcoxon effect size,", "comment"),
    ("# valid across differently-normalized expression scales", "comment"),
    ("meta <- meta_cor_DL(summary_df$r_rb, summary_df$n, summary_df$cohort)", "key"),
    ("", "code"),
    ("# >>> Pooled (random-effects): r=-0.59, p=0.22 (I^2=87.1%)", "comment"),
    ("# ==> 2/3 individually significant, NOT a confirmed pooled finding", "comment"),
]
content_slide(
    "12_hcc_ICI_pooled_summary.R — Suggestive Lead, Not Yet Proven",
    code,
    "연구계획서 \"preliminary data\"로 쓸 때 중요: 3개 중 2개 코호트는 개별적으로 유의하지만, "
    "정식 pooled meta-analysis는 비유의(p=0.22, 높은 이질성). \"이미 증명됨\"이 아니라 \"더 큰 "
    "검증 코호트가 필요한 유망한 신호\"로 제시해야 함.",
    result_image=os.path.join(FIG, "FOREST_hcc_ICI_HSF2_pooled.png"),
)

# =====================================================================
# SLIDE 17: SUMMARY
# =====================================================================
add_overview_slide("Summary of Dry-Lab Findings", [
    "- 검증한 가설: 섬유화 조직에서 HSF2↑ → NETosis 유전자 발현↑(양의 상관). 결과: 4개의 독립적인 "
    "human 간·폐 코호트에서 HSF2는 NETosis_core score와 음의 상관(pooled rho=-0.31, p=0.010), "
    "3/4 코호트에서 재현, PTPRC 면역침윤 대조도 통과한 특이적 신호.",
    "",
    "- 기존 공개 ChIP-seq/co-expression 데이터베이스 어디에도 HSF2가 이 유전자들을 활성화한다는 "
    "근거 없음 — 정말 미개척 영역.",
    "",
    "- 3개의 독립적인 human HCC 면역항암제 코호트에서 종양 HSF2가 낮을수록 치료 반응이 좋아지는 "
    "경향(3개 중 2개는 개별 유의; pooled 결과는 아직 불확실, 추가 데이터 필요).",
    "",
    "- 제안하는 재해석: HSF2는 N2 극성화와 NETosis를 함께 켜는 \"총사령관\"이 아니라, 오히려 둘을 "
    "\"분리(decouple)\"시켜서 — 더 넓은 N2/면역억제 프로그램은 진행되게 두면서 NETosis 쪽에만 "
    "브레이크를 거는 조절자일 수 있음.",
])

# =====================================================================
# SLIDE 18: ASK
# =====================================================================
add_overview_slide("What Would Help Next: The In Vivo RNA-seq Gene List", [
    "- 연구실의 첫 실험(PHMG 6 mg/kg in vivo → 호중구 RNA-seq → 상승 유전자군의 1순위 상위 TF가 "
    "HSF2)이 바로 위 \"decoupling\" 가설을 직접 검증할 수 있는 빠진 조각입니다.",
    "",
    "- 요청: 그 in vivo 실험에서 상승(upregulated)한 유전자 리스트(또는 TF-enrichment 분석에 "
    "입력한 유전자셋).",
    "",
    "- 만약 PADI4/ELANE/MPO/CAMP/DEFA4가 그 상승 유전자 목록에 포함되어 있었다면 → in vivo에서는 "
    "HSF2와 NETosis 유전자가 같이 움직였다는 뜻 → human 조직 결과와 상충, 조정이 필요.",
    "",
    "- 만약 그 유전자들이 상승 목록에 없었다면(변화 없음 또는 감소) → 외부 공개 데이터뿐 아니라 "
    "우리 연구실 자체 마우스 데이터로도 decoupling 가설이 직접 뒷받침됨.",
])

# =====================================================================
# SLIDE 19: NEXT WEEK'S WET-LAB SUGGESTION
# =====================================================================
add_overview_slide("Suggestion for Next Week's Neutrophil Time-Course Experiment", [
    "- 계획: 분리된 호중구에서 PHMG 농도/시간별 실험(기존 3 ug/mL, 6시간 단일 Western blot 결과에서 "
    "확장).",
    "",
    "- 제안하는 추가사항: 단일 시점이 아닌 전체 time course(1h/2h/4h/6h)로 (a) HSF2 단백질/mRNA와 "
    "(b) NETosis 핵심 mRNA(PADI4, ELANE, MPO, CAMP, DEFA4, qPCR)를 함께 추적하고, 표준 N1(iNOS, "
    "TNF-α, ICAM-1)·N2(Arg-1, CD206, IL-10) 마커도 같이, 여기에 기능적 NETosis readout(세포외 "
    "DNA/SYTOX assay) 하나 추가.",
    "",
    "- Dry-lab 데이터에 근거한 예측: HSF2가 오를 때 PADI4/ELANE/MPO/CAMP/DEFA4는 같이 오르지 "
    "않고, Arg-1/CD206은 HSF2와 함께 오를 것 — 독성 스트레스 하에서 N2 극성화와 NETosis 유전자 "
    "유도가 decoupling됨.",
    "",
    "- 어느 쪽 결과가 나오든(확인되든 반박되든) 그 자체로 의미 있는, 보고 가능한 발견입니다.",
])

prs.save(OUT)
print(f"Presentation build complete: {OUT}")
print(f"Total slides: {len(prs.slides)}")
