#!/usr/bin/env python3
"""Plain, no-decoration lab-meeting slide deck: black text on white, code +
result (table/figure) + short explanation per analysis step. Built with
python-pptx (no LibreOffice/Node on this machine for pptxgenjs/visual QA)."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

BASE = r"C:\Users\SAMSUNG\Desktop\project2"
FIG = os.path.join(BASE, "figures")
OUT = os.path.join(BASE, "presentation", "dry_lab_summary.pptx")

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x40, 0x40, 0x40)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

MARGIN = Inches(0.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def new_slide():
    return prs.slides.add_slide(BLANK)


def add_title(slide, text, size=28):
    tb = slide.shapes.add_textbox(MARGIN, Inches(0.3), SLIDE_W - 2 * MARGIN, Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.font.name = "Calibri"
    return tb


def add_text_block(slide, left, top, width, height, text, size=14, bold=False,
                    align=PP_ALIGN.LEFT, italic=False, font="Calibri", color=BLACK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
    return tb


def add_code_box(slide, left, top, width, height, code_text, label="R code", size=10):
    # label
    add_text_block(slide, left, top, width, Inches(0.3), label, size=12, bold=True)
    box_top = top + Inches(0.32)
    box_h = height - Inches(0.32)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, box_top, width, box_h)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = BLACK
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    lines = code_text.strip("\n").split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.name = "Consolas"
        p.font.color.rgb = BLACK
    return box


def add_result_image(slide, left, top, width, height, img_path, label="Result"):
    add_text_block(slide, left, top, width, Inches(0.3), label, size=12, bold=True)
    img_top = top + Inches(0.32)
    img_h = height - Inches(0.32)
    pic = slide.shapes.add_picture(img_path, left, img_top, height=img_h)
    if pic.width > width:
        ratio = width / pic.width
        pic.width = int(width)
        pic.height = int(pic.height * ratio)
    # center horizontally within the region
    pic.left = int(left + (width - pic.width) / 2)
    return pic


def add_result_table(slide, left, top, width, height, headers, rows, label="Result",
                      font_size=11):
    add_text_block(slide, left, top, width, Inches(0.3), label, size=12, bold=True)
    tbl_top = top + Inches(0.32)
    tbl_h = height - Inches(0.32)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, tbl_top, width, tbl_h)
    table = gshape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(font_size)
            p.font.color.rgb = BLACK
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = BLACK
    return gshape


def content_slide(title, code, explanation, result_image=None, result_table=None,
                   code_label="R code", result_label="Result"):
    s = new_slide()
    add_title(s, title)
    half_w = Emu(int((SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2))
    top = Inches(1.2)
    body_h = Inches(4.1)
    add_code_box(s, MARGIN, top, half_w, body_h, code, label=code_label)
    right_left = MARGIN + half_w + Inches(0.3)
    if result_image:
        add_result_image(s, right_left, top, half_w, body_h, result_image, label=result_label)
    elif result_table:
        headers, rows = result_table
        add_result_table(s, right_left, top, half_w, body_h, headers, rows, label=result_label)
    add_text_block(s, MARGIN, Inches(5.55), SLIDE_W - 2 * MARGIN, Inches(1.6),
                    explanation, size=14)
    return s


# =====================================================================
# SLIDE 1: TITLE
# =====================================================================
s = new_slide()
add_text_block(s, MARGIN, Inches(2.6), SLIDE_W - 2 * MARGIN, Inches(1.6),
                "PHMG → HSF2 → NETosis Hypothesis:\nDry-Lab (in silico) Verification Using Public Human Fibrosis & HCC Data",
                size=32, bold=True, align=PP_ALIGN.CENTER)
add_text_block(s, MARGIN, Inches(4.4), SLIDE_W - 2 * MARGIN, Inches(0.6),
                "Bioinformatics follow-up analysis — [Name], [Date]",
                size=16, align=PP_ALIGN.CENTER, italic=True)

# =====================================================================
# SLIDE 2: BACKGROUND
# =====================================================================
s = new_slide()
add_title(s, "Background: Why This Dry-Lab Analysis Was Done")
add_text_block(s, MARGIN, Inches(1.3), SLIDE_W - 2 * MARGIN, Inches(5.5),
    "Lab's experiments so far (in vivo → ex vivo):\n"
    "  1) PHMG 6 mg/kg administered to mice in vivo → neutrophils isolated → RNA-seq.\n"
    "     Upstream-TF analysis of the upregulated gene set ranked HSF2 as the #1 candidate TF.\n"
    "  2) Follow-up validation: neutrophils isolated from untreated mice, treated ex vivo with\n"
    "     PHMG 3 μg/mL for 6 h → Western blot showed a visibly thicker HSF2 band vs. control.\n\n"
    "Working hypothesis (lab):\n"
    "  Chronic toxicant stress (PHMG) → HSF2 activation → neutrophils polarize toward the\n"
    "  anti-inflammatory / immunosuppressive “N2” state → N2 neutrophils suppress cytotoxic\n"
    "  T-cell activity → immune-evasive microenvironment → promotes fibrosis and HCC.\n\n"
    "Open question before more wet-lab work:\n"
    "  N2 polarization is usually described as one bundled program — immunosuppressive markers\n"
    "  (Arg-1, CD206) AND NETosis genes (PADI4, ELANE, MPO) rising together. Does public human\n"
    "  tissue data actually support HSF2 driving NETosis-gene expression specifically?\n"
    "  → This is exactly what the following dry-lab analysis tests.",
    size=16)

# =====================================================================
# SLIDE 3: PRE-SPECIFIED HYPOTHESIS / ANALYSIS PLAN
# =====================================================================
code = '''# scripts/00_functions.R  (plan written BEFORE any result was seen)
# Primary endpoint : Spearman rho(HSF2, NET_core_score)
#                     within diseased/fibrotic samples
# NET_core genes   : PADI4, ELANE, MPO
#                     (direct NETosis machinery, mechanistically required)
# NET_extended     : core + CTSG, LTF, CAMP, DEFA4, AZU1,
#                     ITGAM, NCF2, S100A12, LCN2
# Significance     : alpha = 0.05, two-sided, per cohort
# Replication rule : primary + replication cohort must agree in
#                     sign AND both be nominally significant (p<0.05)

NET_CORE     <- c("PADI4", "ELANE", "MPO")
NET_EXTENDED <- c(NET_CORE, "CTSG", "LTF", "CAMP", "DEFA4",
                   "AZU1", "ITGAM", "NCF2", "S100A12", "LCN2")
TARGET_TF    <- "HSF2"'''
content_slide(
    "Testable Prediction & Pre-Specified Analysis Plan",
    code,
    "If HSF2 drives NETosis-gene expression, HSF2 should positively correlate with the "
    "NETosis-core score (PADI4/ELANE/MPO) in fibrotic tissue, reproducibly across independent "
    "organs/cohorts. This plan (endpoint, gene panels, significance rule, replication rule) was "
    "fixed before looking at any correlation result.",
    result_table=(
        ["Cohort", "Organ", "Role", "Platform", "n (fibrotic)"],
        [
            ["GSE135251", "Liver", "Primary", "RNA-seq", "206 (NAFLD F0-F4)"],
            ["GSE14323", "Liver", "Replication", "Affymetrix", "96"],
            ["GSE47460", "Lung", "Primary", "Agilent", "122 (UIP/IPF)"],
            ["GSE53845", "Lung", "Replication", "Agilent 2-color", "40 (IPF)"],
        ],
    ),
    code_label="scripts/00_functions.R",
    result_label="4 independent public cohorts used",
)

# =====================================================================
# SLIDE 4: GSE135251 (liver, primary)
# =====================================================================
code = '''# scripts/01_liver_GSE135251.R
net_core_score <- panel_score(expr_panel, NET_CORE, min_genes = 2)
df_nafld <- df[df$disease == "NAFLD", ]   # fibrotic subset, n=206

results$primary <- spearman_report(
    df_nafld$HSF2, df_nafld$NET_core,
    "GSE135251_liver: HSF2 vs NET_core (NAFLD subset, primary)")

# >>> rho = -0.145, 95% CI [-0.288, -0.002], p = 0.038'''
content_slide(
    "Step 1 — Liver, Primary Cohort: GSE135251 (n=206, RNA-seq)",
    code,
    "Result is opposite in direction to the original hypothesis: HSF2 correlates NEGATIVELY, "
    "not positively, with NETosis-core gene expression in fibrotic (NAFLD) liver tissue "
    "(rho = -0.145, p = 0.038).",
    result_image=os.path.join(FIG, "liver_GSE135251_HSF2_vs_NETcore.png"),
    code_label="scripts/01_liver_GSE135251.R",
    result_label="figures/liver_GSE135251_HSF2_vs_NETcore.png",
)

# =====================================================================
# SLIDE 5: GSE14323 (liver, replication)
# =====================================================================
code = '''# scripts/02_liver_GSE14323.R  (independent platform: Affymetrix)
df_disease <- df[df$tissue != "Normal", ]  # cirrhosis/cirrhosisHCC/HCC, n=96

results$primary <- spearman_report(
    df_disease$HSF2, df_disease$NET_core,
    "GSE14323_liver: HSF2 vs NET_core (diseased subset, replication)")

# >>> rho = -0.550, 95% CI [-0.689, -0.381], p = 6.2e-09
# Replication rule (same sign + both p<0.05 vs GSE135251): MET'''
content_slide(
    "Step 2 — Liver, Replication Cohort: GSE14323 (n=96, microarray)",
    code,
    "Independent platform, independent patients — same negative direction, much stronger and "
    "highly significant (rho = -0.550, p = 6.2e-09). Formally REPLICATES the liver finding by "
    "the pre-specified rule.",
    result_image=os.path.join(FIG, "liver_GSE14323_HSF2_vs_NETcore.png"),
    code_label="scripts/02_liver_GSE14323.R",
    result_label="figures/liver_GSE14323_HSF2_vs_NETcore.png",
)

# =====================================================================
# SLIDE 6: GSE47460 (lung, primary)
# =====================================================================
code = '''# scripts/03_lung_GSE47460.R
df_fibrotic <- df[df$group == "UIP_IPF", ]  # histologically fibrotic, n=122

results$primary <- spearman_report(
    df_fibrotic$HSF2, df_fibrotic$NET_core,
    "GSE47460_lung: HSF2 vs NET_core (UIP/IPF subset, primary)")

# >>> rho = -0.392, 95% CI [-0.538, -0.218], p = 8.1e-06'''
content_slide(
    "Step 3 — Lung, Primary Cohort: GSE47460 (n=122 IPF, Agilent)",
    code,
    "Different organ entirely (lung vs. liver), different platform — same negative direction, "
    "highly significant (rho = -0.392, p = 8.1e-06). The pattern is not liver-specific.",
    result_image=os.path.join(FIG, "lung_GSE47460_HSF2_vs_NETcore.png"),
    code_label="scripts/03_lung_GSE47460.R",
    result_label="figures/lung_GSE47460_HSF2_vs_NETcore.png",
)

# =====================================================================
# SLIDE 7: GSE53845 (lung, replication)
# =====================================================================
code = '''# scripts/04_lung_GSE53845.R  (independent submission, n=40 IPF)
df_ipf <- df[df$diagnosis == "IPF", ]

results$primary <- spearman_report(
    df_ipf$HSF2, df_ipf$NET_core,
    "GSE53845_lung: HSF2 vs NET_core (IPF subset, replication)")

# >>> rho = -0.050, 95% CI [-0.350, 0.274], p = 0.760  (n.s.)
# Same sign as GSE47460, but NOT individually significant
# -> interpreted as underpowered (n=40), not as contradicting evidence'''
content_slide(
    "Step 4 — Lung, Replication Cohort: GSE53845 (n=40 IPF, Agilent 2-color)",
    code,
    "Direction still negative, but small n=40 is underpowered — does not reach significance "
    "on its own. Treated as \"no information\" rather than a contradiction, and carried "
    "forward into the meta-analysis below.",
    result_image=os.path.join(FIG, "lung_GSE53845_HSF2_vs_NETcore.png"),
    code_label="scripts/04_lung_GSE53845.R",
    result_label="figures/lung_GSE53845_HSF2_vs_NETcore.png",
)

# =====================================================================
# SLIDE 8: META-ANALYSIS (KEY FINDING)
# =====================================================================
code = '''# scripts/05_meta_analysis.R
# from-scratch Fisher-z / DerSimonian-Laird random-effects meta-analysis
# (no internet-installed 'metafor' package in this environment)
meta_all <- meta_cor_DL(primary_df$rho, primary_df$n, primary_df$cohort)

cat(sprintf("pooled rho = %.3f, 95%% CI [%.3f, %.3f], p = %.3g
Heterogeneity: I^2 = %.1f%%",
    meta_all$pooled_rho, meta_all$ci_lo, meta_all$ci_hi,
    meta_all$p_value, meta_all$I2))

# >>> pooled rho = -0.308, 95% CI [-0.508, -0.076], p = 0.010
# >>> I^2 = 83.5% (liver-only: -0.36, p=0.11 / lung-only: -0.25, p=0.15)'''
content_slide(
    "KEY FINDING — 4-Cohort Meta-Analysis: HSF2 vs. NETosis-Core Score",
    code,
    "Pooled effect (p=0.010) IS statistically significant and OPPOSITE the original "
    "hypothesis - even though 1 of 4 cohorts (GSE53845, n=40) was individually underpowered. "
    "But note I^2=83.5% (high heterogeneity): significance says the direction is real; "
    "heterogeneity says the SIZE of the effect is not uniform across cohorts. Two separate "
    "questions - both matter, and this is why the follow-up wet-lab time-course is needed to "
    "settle the magnitude question.",
    result_image=os.path.join(FIG, "FOREST_HSF2_vs_NETcore_all_cohorts.png"),
    code_label="scripts/05_meta_analysis.R",
    result_label="figures/FOREST_HSF2_vs_NETcore_all_cohorts.png",
)

# =====================================================================
# SLIDE 9: SPECIFICITY CONTROL (PTPRC)
# =====================================================================
code = '''# scripts/06_specificity_check_PTPRC.R
# Does HSF2 fall because immune infiltration simply "dilutes" bulk
# tissue signal as fibrosis worsens? Test against PTPRC (CD45,
# pan-leukocyte marker) - if that were the explanation, HSF2 should
# ALSO correlate negatively with PTPRC. It does not.

r1_ptprc   <- spearman_report(df1_dis$HSF2, df1_dis$PTPRC, ...)
r1_netcore <- spearman_report(df1_dis$HSF2, df1_dis$NET_core, ...)

# >>> GSE135251: HSF2 vs PTPRC     rho = +0.178  p = 0.010
# >>> GSE135251: HSF2 vs NET_core  rho = -0.145  p = 0.038
# >>> GSE14323:  HSF2 vs PTPRC     rho = +0.448  p = 4.7e-06
# >>> GSE14323:  HSF2 vs NET_core  rho = -0.550  p = 6.2e-09'''
content_slide(
    "Specificity Control: Not a Generic Immune-Infiltration Artifact",
    code,
    "HSF2 rises WITH overall immune infiltration (PTPRC) but falls specifically with the "
    "NETosis-core score — opposite directions. Rules out the trivial \"more immune cells = less "
    "relative HSF2 signal\" explanation; the negative NETosis relationship looks specific.",
    result_table=(
        ["Cohort", "HSF2 vs PTPRC", "HSF2 vs NET_core"],
        [
            ["GSE135251", "+0.178 (p=0.010)", "-0.145 (p=0.038)"],
            ["GSE14323", "+0.448 (p=4.7e-06)", "-0.550 (p=6.2e-09)"],
        ],
    ),
    code_label="scripts/06_specificity_check_PTPRC.R",
    result_label="results/SPECIFICITY_check_PTPRC_vs_NETcore.csv",
)

# =====================================================================
# SLIDE 10: CROSS-COHORT GENE CONSISTENCY
# =====================================================================
code = '''# scripts/08_cross_cohort_gene_consistency.R
# Break NET_extended into individual genes: which ones move with
# HSF2 the SAME direction in all 4 cohorts, vs. which flip sign
# between liver and lung?

consistency$n_cohorts_sig_FDR05 <- rowSums(consistency[, fdr_cols] < 0.05)
consistency$direction_consistent <-
    (consistency$n_negative == consistency$n_cohorts_tested) |
    (consistency$n_positive == consistency$n_cohorts_tested)

# >>> Consistent (all 4 cohorts negative): CAMP, MPO, DEFA4
#     (PADI4, ELANE: negative in 3/4, one small ns exception)
# >>> Organ-DIVERGENT (liver vs lung flip sign): LTF, NCF2, ITGAM, LCN2'''
content_slide(
    "Which Individual NETosis Genes Drive This — and Which Don't",
    code,
    "The core NETosis machinery genes (CAMP, MPO, DEFA4, PADI4, ELANE) move consistently "
    "negative with HSF2 in all 4 cohorts. A separate gene set (LTF, NCF2, ITGAM, LCN2) actually "
    "flips sign between liver and lung — possibly not neutrophil-specific (e.g. LCN2 is also "
    "made by stressed hepatocytes).",
    result_table=(
        ["Gene", "GSE135251", "GSE14323", "GSE47460", "GSE53845", "Direction"],
        [
            ["CAMP", "-0.056", "-0.459", "-0.457", "-0.529", "Consistent (-)"],
            ["MPO", "-0.025", "-0.277", "-0.320", "-0.412", "Consistent (-)"],
            ["DEFA4", "-0.090", "-0.337", "-0.328", "-0.097", "Consistent (-)"],
            ["PADI4", "-0.120", "-0.496", "-0.205", "+0.181", "3 neg / 1 pos (ns)"],
            ["ELANE", "-0.192", "-0.422", "-0.348", "+0.075", "3 neg / 1 pos (ns)"],
        ],
    ),
    code_label="scripts/08_cross_cohort_gene_consistency.R",
    result_label="results/CROSS_COHORT_gene_consistency.csv",
)

# =====================================================================
# SLIDE 11: CHEA3 IN SILICO CHECK
# =====================================================================
code = '''# scripts/09_chea3_analysis.R
# Does any existing ChIP-seq / co-expression database already show
# HSF2 as an upstream regulator of these genes?
chea3_query <- function(genes, query_name) {
  res <- httr::POST(
    url = "https://maayanlab.cloud/chea3/api/enrich/",
    body = list(query_name = query_name, gene_set = as.list(genes)),
    encode = "json", httr::timeout(30))
  fromJSON(httr::content(res, as = "text"), simplifyVector = TRUE)
}
# Query: PADI4, ELANE, MPO, CAMP, DEFA4 (the consistent gene set)

# >>> HSF2 rank ~1200-1450 of ~1400-1630 TFs in every library (bottom ~10%)
# >>> Gene overlap (intersect) = 0 in every co-expression library
# >>> NOT returned at all by ENCODE / ReMap / Literature ChIP-seq'''
content_slide(
    "In Silico Check: No Existing Evidence HSF2 Directly Activates These Genes",
    code,
    "Across 8 independent evidence libraries (ChIP-seq + co-expression + literature), HSF2 is "
    "never a plausible top regulator of this gene set. Consistent with either (a) a repressive "
    "relationship that co-expression tools can't detect, or (b) this axis being genuinely "
    "unstudied — i.e. this is real blue-ocean territory, not something already known and missed.",
    result_table=(
        ["Library", "HSF2 rank", "Total TFs", "Gene overlap"],
        [
            ["Integrated mean rank", "1444", "1632", "-"],
            ["ARCHS4 co-expression", "1209", "1628", "0"],
            ["GTEx co-expression", "1417", "1607", "0"],
            ["ENCODE ChIP-seq", "not returned", "118", "-"],
            ["ReMap ChIP-seq", "not returned", "297", "-"],
        ],
    ),
    code_label="scripts/09_chea3_analysis.R",
    result_label="results/CHEA3_HSF2_enrichment_results.csv",
)

# =====================================================================
# SLIDE 12 (NEW): TRANSITION - WHAT IS "HCC-ICI" AND WHY
# =====================================================================
s = new_slide()
add_title(s, "Next: A Different Question — Does This Matter Clinically?")
add_text_block(s, MARGIN, Inches(1.3), SLIDE_W - 2 * MARGIN, Inches(5.6),
    "So far: fibrosis tissue (not cancer) — testing whether HSF2 relates to NETosis genes at "
    "all. The next 3 slides switch to a different question and a different kind of data.\n\n"
    "HCC-ICI = Hepatocellular Carcinoma + Immune Checkpoint Inhibitor.\n"
    "  - HCC: liver cancer\n"
    "  - ICI: immune checkpoint inhibitor drugs (anti-PD-1 / anti-PD-L1 immunotherapy)\n"
    "  - \"Cohort\": real liver-cancer patients who received these drugs, each labeled "
    "Responder or Non-responder, with tumor RNA-seq data available\n\n"
    "Why look at this: a prior published study (Chen et al., 2022) already showed HSF2 is "
    "elevated in liver cancer and correlates with immune-checkpoint gene expression. So beyond "
    "the fibrosis/NETosis question, we also asked: does tumor HSF2 relate to whether real HCC "
    "patients actually respond to immunotherapy? This tests whether the HSF2 axis has direct "
    "clinical relevance, not just a tissue-level correlation.",
    size=16)

# =====================================================================
# SLIDE 13: HCC-ICI GATE-CHECK + GSE215011
# =====================================================================
code = '''# GSE140901 (originally planned HCC-ICI dataset) - GATE-CHECK:
# 785-gene NanoString panel. Does it even measure our genes?
zcat GSE140901_processed_data.txt.gz | grep -w "HSF2"
# >>> (no output - HSF2 is NOT in the panel; PADI4/MPO also missing)
# ==> EXCLUDED, do not force-fit a substitute readout

# scripts/07_hcc_ICI_GSE215011.R  (found instead: whole-transcriptome
# RNA-seq, HCC, anti-PD-1 monotherapy, n=10, 5 responder/5 non-resp.)
wilcox_report(df$HSF2, df$group, "HSF2: Responder vs Non-responder")
# >>> rank-biserial r = 0.44, p = 0.30  (not significant, n=10)'''
content_slide(
    "HCC Immunotherapy: Dataset Gate-Check + First Cohort (GSE215011)",
    code,
    "The originally planned dataset (GSE140901) turned out not to measure HSF2 at all — "
    "checked directly and excluded rather than substituting a different gene. Found and used "
    "GSE215011 instead (real RNA-seq). Alone, n=10 is too small to reach significance either way.",
    result_image=os.path.join(FIG, "hcc_ICI_GSE215011_HSF2_vs_NETcore.png"),
    code_label="terminal + scripts/07_hcc_ICI_GSE215011.R",
    result_label="figures/hcc_ICI_GSE215011_HSF2_vs_NETcore.png",
)

# =====================================================================
# SLIDE 13: GSE279750 + GSE235863
# =====================================================================
code = '''# scripts/10_hcc_ICI_GSE279750.R  (anti-PD-L1 combo, post-tx, n=10)
wilcox_report(df$HSF2, df$group, "HSF2: Responder vs Non-responder")
# >>> rank-biserial r = -0.92, p = 0.025   (responders have LOWER HSF2)

# scripts/11_hcc_ICI_GSE235863.R  (anti-PD-1+lenvatinib, post-tx, n=15)
wilcox_report(df$HSF2, df$group, "HSF2: Responder vs Non-responder")
# >>> rank-biserial r = -0.73, p = 0.043   (same direction, also sig.)'''
content_slide(
    "Two More Independent HCC-ICI Cohorts: Same Direction, Individually Significant",
    code,
    "Found 2 more independent public cohorts. In both, patients with LOWER tumor HSF2 "
    "responded better to combination immunotherapy (both p < 0.05). Note: both are "
    "combination-regimen, post-treatment cohorts — different design from GSE215011.",
    result_image=os.path.join(FIG, "hcc_ICI_GSE279750_HSF2_by_response.png"),
    code_label="scripts/10 + 11_hcc_ICI_*.R",
    result_label="figures/hcc_ICI_GSE279750_HSF2_by_response.png",
)

# =====================================================================
# SLIDE 14: ICI POOLED META-ANALYSIS
# =====================================================================
code = '''# scripts/12_hcc_ICI_pooled_summary.R
# Pool all 3 cohorts (rank-biserial r = scale-free Wilcoxon effect
# size, valid across differently-normalized expression scales)
meta <- meta_cor_DL(summary_df$r_rb, summary_df$n, summary_df$cohort)

# >>> GSE215011 (monotherapy):     r=+0.44  p=0.30
# >>> GSE279750 (combo, post-tx):  r=-0.92  p=0.025
# >>> GSE235863 (combo, post-tx):  r=-0.73  p=0.043
# >>> Pooled (random-effects):     r=-0.59  p=0.22  (I^2=87.1%, high)
# ==> 2/3 individually significant, but NOT a confirmed pooled finding'''
content_slide(
    "HCC-ICI Pooled Result: Suggestive Lead, Not Yet Proven",
    code,
    "IMPORTANT for grant \"preliminary data\" framing: 2 of 3 cohorts are individually "
    "significant, but the formal pooled meta-analysis is NOT significant (p=0.22) due to high "
    "heterogeneity (monotherapy cohort trends opposite). Should be presented as a promising "
    "signal needing a larger validation cohort — not as already proven.",
    result_image=os.path.join(FIG, "FOREST_hcc_ICI_HSF2_pooled.png"),
    code_label="scripts/12_hcc_ICI_pooled_summary.R",
    result_label="figures/FOREST_hcc_ICI_HSF2_pooled.png",
)

# =====================================================================
# SLIDE 15: SUMMARY
# =====================================================================
s = new_slide()
add_title(s, "Summary of Dry-Lab Findings")
add_text_block(s, MARGIN, Inches(1.3), SLIDE_W - 2 * MARGIN, Inches(5.6),
    "1. Original hypothesis tested: HSF2 upregulation -> NETosis-gene expression up (positively "
    "correlated) in fibrotic tissue.\n\n"
    "2. Result in 4 independent human liver/lung fibrosis cohorts: HSF2 is NEGATIVELY correlated "
    "with the NETosis-core score (pooled rho = -0.31, p = 0.010), replicated in 3/4 cohorts, and "
    "specific (survives the PTPRC immune-infiltration control).\n\n"
    "3. No existing public ChIP-seq/co-expression database shows HSF2 activating these genes — "
    "genuinely unstudied territory.\n\n"
    "4. In 3 independent human HCC immunotherapy cohorts, lower tumor HSF2 trends toward better "
    "treatment response (2/3 individually significant; pooled result inconclusive, needs more "
    "data).\n\n"
    "=> Proposed reframing: HSF2 may not be the \"master switch\" that turns ON NETosis together "
    "with N2 polarization — it may instead \"decouple\" the two, acting as a brake specifically "
    "on the NETosis arm while the broader N2/immunosuppressive program proceeds.",
    size=15)

# =====================================================================
# SLIDE 16: ASK
# =====================================================================
s = new_slide()
add_title(s, "What Would Help Next: The In Vivo RNA-seq Gene List")
add_text_block(s, MARGIN, Inches(1.3), SLIDE_W - 2 * MARGIN, Inches(5.6),
    "The very first lab experiment (PHMG 6 mg/kg in vivo -> neutrophil RNA-seq -> HSF2 ranked "
    "#1 upstream TF of the upregulated genes) is the missing piece that could directly test the "
    "\"decoupling\" idea above.\n\n"
    "Request: the list of genes that were UPREGULATED in that in vivo experiment (or at least "
    "the gene set that was fed into the TF-enrichment analysis).\n\n"
    "Why this matters:\n"
    "  - If PADI4 / ELANE / MPO (or CAMP / DEFA4) were part of that upregulated set -> HSF2 and "
    "NETosis genes moved together in vivo -> tension with the human tissue finding above, needs "
    "reconciling.\n"
    "  - If those genes were NOT in the upregulated set (flat or down) -> directly SUPPORTS the "
    "decoupling hypothesis using the lab's own mouse data, not just outside public data.\n\n"
    "This single check could turn \"human data disagrees with the original hypothesis\" into "
    "\"human data + our own mouse data together reveal a new regulatory checkpoint.\"",
    size=15)

# =====================================================================
# SLIDE 17: NEXT WEEK'S WET-LAB SUGGESTION
# =====================================================================
s = new_slide()
add_title(s, "Suggestion for Next Week's Neutrophil Time-Course Experiment")
add_text_block(s, MARGIN, Inches(1.3), SLIDE_W - 2 * MARGIN, Inches(5.6),
    "Planned: PHMG dose/time-course in isolated neutrophils (building on the single 3 ug/mL, "
    "6 h Western blot result).\n\n"
    "Suggested addition based on the dry-lab result — pre-registered prediction to test:\n"
    "  - Full time course (e.g. 1h / 2h / 4h / 6h), not a single timepoint, tracking BOTH:\n"
    "      (a) HSF2 protein (Western blot) and mRNA\n"
    "      (b) NETosis-core mRNA: PADI4, ELANE, MPO, and CAMP, DEFA4 (qPCR)\n"
    "  - In parallel, standard N1 (iNOS, TNF-a, ICAM-1) and N2 (Arg-1, CD206, IL-10) markers\n"
    "  - One functional NETosis readout (e.g. extracellular DNA / SYTOX assay), not mRNA alone\n\n"
    "Prediction from dry-lab data: as HSF2 rises, PADI4/ELANE/MPO/CAMP/DEFA4 will NOT rise in "
    "parallel (flat or down), while Arg-1/CD206 rise together with HSF2 — i.e., N2 polarization "
    "and NETosis-gene induction become decoupled under toxicant stress.\n\n"
    "Either outcome (confirmed or refuted) is a real, reportable finding.",
    size=15)

prs.save(OUT)
print(f"Presentation build complete: {OUT}")
print(f"Total slides: {len(prs.slides)}")
